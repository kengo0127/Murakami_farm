from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── カラー定義 ──────────────────────────────────────────────
DARK_GREEN  = RGBColor(0x1a, 0x3d, 0x17)
MED_GREEN   = RGBColor(0x2e, 0x5c, 0x28)
LIGHT_GREEN = RGBColor(0xe8, 0xf5, 0xe9)
ACCENT      = RGBColor(0x4c, 0xaf, 0x50)
WHITE       = RGBColor(0xff, 0xff, 0xff)
DARK_GRAY   = RGBColor(0x22, 0x22, 0x22)
MID_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xf5, 0xf5, 0xf5)
ORANGE      = RGBColor(0xff, 0x98, 0x00)
RED         = RGBColor(0xf4, 0x43, 0x36)
SLATE       = RGBColor(0x90, 0xa4, 0xae)

FONT   = "Meiryo"
W      = Inches(13.33)   # スライド幅
H      = Inches(7.5)     # スライド高さ


# ── ヘルパー関数 ──────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # 完全ブランク
    return prs.slides.add_slide(layout)

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, wrap=True, font=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox

def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=DARK_GRAY, indent_items=None):
    """items: list of str  |  indent_items: set of indices that are sub-bullets"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    indent_items = indent_items or set()
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        is_sub = i in indent_items
        p.level = 1 if is_sub else 0
        run = p.add_run()
        run.text = item
        run.font.name  = FONT
        run.font.size  = Pt(font_size - 2 if is_sub else font_size)
        run.font.color.rgb = MID_GRAY if is_sub else color
    return txBox

def header_bar(slide, title, subtitle=None):
    """ダークグリーンのヘッダーバー"""
    bar_h = Inches(1.2)
    add_rect(slide, 0, 0, W, bar_h, DARK_GREEN)
    # ロゴアイコン
    add_text(slide, "🌿", Inches(0.2), Inches(0.15), Inches(0.6), Inches(0.9),
             font_size=36, color=WHITE)
    # タイトル
    t_left = Inches(0.85)
    if subtitle:
        add_text(slide, title, t_left, Inches(0.1), W - t_left - Inches(0.2), Inches(0.6),
                 font_size=22, bold=True, color=WHITE)
        add_text(slide, subtitle, t_left, Inches(0.65), W - t_left - Inches(0.2), Inches(0.45),
                 font_size=13, color=RGBColor(0xa5, 0xd6, 0xa7))
    else:
        add_text(slide, title, t_left, Inches(0.2), W - t_left - Inches(0.2), Inches(0.8),
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def section_badge(slide, text, left, top, bg=MED_GREEN):
    """小さなセクションラベル"""
    w, h = Inches(1.8), Inches(0.36)
    add_rect(slide, left, top, w, h, bg)
    add_text(slide, text, left, top, w, h,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def step_box(slide, num, title, body, left, top, width, height, accent=MED_GREEN):
    """番号付きステップボックス"""
    # 番号円
    cw = Inches(0.55)
    add_rect(slide, left, top + (height - cw) / 2, cw, cw, accent)
    add_text(slide, str(num), left, top + (height - cw) / 2, cw, cw,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # テキスト
    add_text(slide, title, left + cw + Inches(0.1), top, width - cw - Inches(0.1), Inches(0.35),
             font_size=14, bold=True, color=MED_GREEN)
    add_text(slide, body, left + cw + Inches(0.1), top + Inches(0.33), width - cw - Inches(0.1),
             height - Inches(0.33), font_size=12, color=MID_GRAY)

def status_legend(slide, left, top):
    """ステータス凡例"""
    items = [
        (ACCENT,  "完了"),
        (ORANGE,  "作業中"),
        (RED,     "遅延（時刻超過・未完了）"),
        (SLATE,   "未着手"),
    ]
    cx = left
    for color, label in items:
        add_rect(slide, cx, top, Inches(0.22), Inches(0.22), color)
        add_text(slide, label, cx + Inches(0.28), top - Inches(0.02),
                 Inches(1.6), Inches(0.28), font_size=12, color=DARK_GRAY)
        cx += Inches(2.1)


# ════════════════════════════════════════════════════════════
# 管理者向けガイド（5枚）
# ════════════════════════════════════════════════════════════

def make_admin_ppt(path):
    prs = new_prs()

    # ── スライド1：タイトル ─────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, DARK_GREEN)
    # 白い中央カード
    cw, ch = Inches(9), Inches(4.6)
    cx = (W - cw) / 2
    cy = (H - ch) / 2
    add_rect(sl, cx, cy, cw, ch, WHITE)
    add_rect(sl, cx, cy, cw, Inches(0.08), ACCENT)   # アクセントライン
    add_text(sl, "🌿  村上農園", cx, cy + Inches(0.3), cw, Inches(0.7),
             font_size=28, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "タスク管理システム", cx, cy + Inches(0.95), cw, Inches(0.7),
             font_size=32, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_rect(sl, cx + Inches(3.2), cy + Inches(1.85), Inches(2.6), Inches(0.05), LIGHT_GREEN)
    add_text(sl, "管理者向け  操作ガイド", cx, cy + Inches(1.95), cw, Inches(0.6),
             font_size=20, bold=False, color=MED_GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "2026年6月", cx, cy + Inches(3.9), cw, Inches(0.4),
             font_size=12, color=SLATE, align=PP_ALIGN.CENTER)

    # ── スライド2：毎朝の操作フロー ────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    header_bar(sl, "毎朝の操作フロー", "管理者が毎日行う3ステップ")

    steps = [
        ("ログイン",
         "スマートフォン or PC で\nhttps://kengo0127.github.io/Murakami_farm/\nにアクセスしてログイン"),
        ("タスク登録",
         "「📝 タスク登録」タブでその日の\n作業タスクをフォームから登録する"),
        ("進捗確認 & エクスポート",
         "「📊 進捗確認」でリアルタイムに確認\n必要に応じて「📥 データ出力」でCSV出力"),
    ]
    box_w = Inches(3.6)
    box_h = Inches(3.8)
    tops  = Inches(1.55)
    for i, (title, body) in enumerate(steps):
        lx = Inches(0.5) + i * (box_w + Inches(0.5))
        add_rect(sl, lx, tops, box_w, box_h, WHITE, ACCENT, 1.5)
        # 番号バッジ
        add_rect(sl, lx + (box_w - Inches(0.7)) / 2, tops - Inches(0.38), Inches(0.7), Inches(0.7), ACCENT)
        add_text(sl, str(i + 1), lx + (box_w - Inches(0.7)) / 2, tops - Inches(0.38),
                 Inches(0.7), Inches(0.7), font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title, lx + Inches(0.2), tops + Inches(0.2), box_w - Inches(0.4), Inches(0.5),
                 font_size=16, bold=True, color=MED_GREEN, align=PP_ALIGN.CENTER)
        add_rect(sl, lx + Inches(0.4), tops + Inches(0.72), box_w - Inches(0.8), Inches(0.03), LIGHT_GREEN)
        add_text(sl, body, lx + Inches(0.2), tops + Inches(0.85), box_w - Inches(0.4), Inches(2.7),
                 font_size=13, color=DARK_GRAY)
        if i < 2:
            add_text(sl, "→", lx + box_w + Inches(0.1), tops + box_h / 2 - Inches(0.25), Inches(0.4),
                     Inches(0.5), font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # ── スライド3：タスク登録 ───────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    header_bar(sl, "タスクの登録方法", "「📝 タスク登録」タブ")

    section_badge(sl, "入力項目", Inches(0.4), Inches(1.45))
    fields = [
        "📅  対象日　　　　　作業日を選択（デフォルト：今日）",
        "🌱  作物　　　　　　スイカ / メロン（入力補完あり）",
        "📝  タスク名　　　　例：洗浄・袋詰め、水やり",
        "👤  担当スタッフ　　チェックボックスで複数選択可",
        "🕐  開始・終了時刻　15分刻みのドロップダウンで選択",
        "💬  メモ（任意）　　補足事項があれば入力",
    ]
    add_bullets(sl, fields, Inches(0.4), Inches(1.9), Inches(6.2), Inches(3.8), font_size=14)

    section_badge(sl, "登録後の操作", Inches(7.0), Inches(1.45), bg=DARK_GREEN)
    ops = [
        "✅  登録済みタスクが一覧に表示される",
        "✏️  「編集」ボタンでフォームに内容を読み込み変更できる",
        "🗑️  「削除」ボタンで完了記録ごと削除される",
        "⚠️  削除すると完了記録も消えるため注意",
    ]
    add_bullets(sl, ops, Inches(7.0), Inches(1.9), Inches(5.9), Inches(3.5), font_size=13)

    tip_top = Inches(5.9)
    add_rect(sl, Inches(0.4), tip_top, Inches(12.5), Inches(1.2), LIGHT_GREEN)
    add_text(sl, "💡  ポイント：スタッフ画面はリアルタイムで更新されます。登録が完了したらスタッフに作業開始を伝えてください。",
             Inches(0.6), tip_top + Inches(0.1), Inches(12.1), Inches(1.0),
             font_size=13, color=MED_GREEN)

    # ── スライド4：進捗確認 ─────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    header_bar(sl, "進捗確認の見方", "「📊 進捗確認」タブ　スタッフの完了状況をリアルタイムで把握")

    # 凡例
    add_text(sl, "ステータスの色分け", Inches(0.4), Inches(1.38), Inches(4), Inches(0.35),
             font_size=13, bold=True, color=MED_GREEN)
    status_legend(sl, Inches(0.4), Inches(1.8))

    # カラム模式図
    crops = [("スイカ", [
        ("07:00〜10:00", "整枝・玉直し", "completed"),
        ("10:30〜12:00", "水やり", "overdue"),
    ]), ("メロン", [
        ("07:00〜10:00", "整枝・玉直し", "completed"),
        ("10:30〜12:00", "水やり", "active"),
    ])]

    col_w = Inches(3.8)
    for ci, (crop, tasks) in enumerate(crops):
        cx = Inches(0.4) + ci * (col_w + Inches(0.3))
        cy = Inches(2.35)
        # カラムヘッダー
        add_rect(sl, cx, cy, col_w, Inches(0.45), DARK_GREEN)
        add_text(sl, crop, cx, cy, col_w, Inches(0.45),
                 font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # タスクカード
        for ti, (time_str, task_name, status) in enumerate(tasks):
            ty = cy + Inches(0.5) + ti * Inches(1.55)
            sc = {
                "completed": ACCENT,
                "active":    ORANGE,
                "overdue":   RED,
                "upcoming":  SLATE,
            }.get(status, SLATE)
            add_rect(sl, cx, ty, col_w, Inches(1.4), WHITE, sc, 3)
            # ステータスバー
            add_rect(sl, cx, ty, Inches(0.08), Inches(1.4), sc)
            add_text(sl, time_str, cx + Inches(0.2), ty + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
                     font_size=11, color=SLATE)
            add_text(sl, task_name, cx + Inches(0.2), ty + Inches(0.38), col_w - Inches(0.3), Inches(0.4),
                     font_size=14, bold=True, color=DARK_GRAY)
            label = {"completed": "✅ 完了", "active": "▶ 作業中", "overdue": "⚠ 遅延", "upcoming": "未着手"}.get(status)
            add_text(sl, label, cx + Inches(0.2), ty + Inches(0.8), col_w - Inches(0.3), Inches(0.35),
                     font_size=12, color=sc, bold=True)

    # 右側説明
    rx = Inches(8.6)
    add_text(sl, "画面の見方", rx, Inches(1.38), Inches(4.5), Inches(0.4),
             font_size=14, bold=True, color=MED_GREEN)
    notes = [
        "・ 作物ごとに縦列で表示",
        "・ タスクは時刻の早い順に並ぶ",
        "・ 赤い横線＝現在時刻",
        "・ スタッフが完了ボタンを押すと\n  自動で画面が更新される",
        "・ 1分ごとに状態を自動チェック",
    ]
    add_bullets(sl, notes, rx, Inches(1.85), Inches(4.5), Inches(3.5), font_size=13)

    # ── スライド5：エクスポート ──────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    header_bar(sl, "データのエクスポート（CSV出力）", "「📥 データ出力」タブ")

    # 操作手順
    section_badge(sl, "操作手順", Inches(0.4), Inches(1.45))
    op_steps = [
        ("①", "「📥 データ出力」タブを選択"),
        ("②", "開始日・終了日を入力して期間を指定"),
        ("③", "「CSVダウンロード」ボタンを押す"),
        ("④", "ダウンロードされた .csv ファイルをExcelで開く"),
    ]
    for i, (num, txt) in enumerate(op_steps):
        ty = Inches(1.95) + i * Inches(0.72)
        add_rect(sl, Inches(0.4), ty, Inches(0.45), Inches(0.45), MED_GREEN)
        add_text(sl, num, Inches(0.4), ty, Inches(0.45), Inches(0.45),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, txt, Inches(0.95), ty + Inches(0.04), Inches(5.5), Inches(0.38),
                 font_size=14, color=DARK_GRAY)

    # 出力項目表
    section_badge(sl, "出力項目（1タスク1行）", Inches(7.2), Inches(1.45), bg=DARK_GREEN)
    cols = [
        ("日付", "作業日"),
        ("スタッフ数", "割り当て人数"),
        ("タスク名 / 作物", ""),
        ("予定開始・終了", ""),
        ("予定所要時間(分)", "終了−開始"),
        ("完了時刻", "実際の時刻"),
        ("所要時間(分)", "開始からの経過"),
        ("予実の差分(分)", "＋遅延 ／ −早期"),
        ("スタッフ名", "全員（読点区切り）"),
    ]
    for i, (col, desc) in enumerate(cols):
        ty = Inches(1.92) + i * Inches(0.5)
        bg = LIGHT_GREEN if i % 2 == 0 else WHITE
        add_rect(sl, Inches(7.2), ty, Inches(2.6), Inches(0.48), bg)
        add_rect(sl, Inches(9.8), ty, Inches(3.3), Inches(0.48), bg)
        add_text(sl, col, Inches(7.25), ty + Inches(0.08), Inches(2.5), Inches(0.34),
                 font_size=12, bold=True, color=MED_GREEN)
        add_text(sl, desc, Inches(9.85), ty + Inches(0.08), Inches(3.2), Inches(0.34),
                 font_size=12, color=MID_GRAY)

    tip_top = Inches(6.1)
    add_rect(sl, Inches(0.4), tip_top, Inches(6.1), Inches(1.0), LIGHT_GREEN)
    add_text(sl, "💡  文字化けしない形式（Excel用BOM付きUTF-8）で出力されます。\n    Excelで直接開いて分析・印刷にご利用ください。",
             Inches(0.6), tip_top + Inches(0.1), Inches(5.9), Inches(0.85), font_size=13, color=MED_GREEN)

    prs.save(path)
    print("admin ppt saved")


# ════════════════════════════════════════════════════════════
# スタッフ向けガイド（3枚）
# ════════════════════════════════════════════════════════════

def make_staff_ppt(path):
    prs = new_prs()

    # ── スライド1：タイトル ─────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, MED_GREEN)
    cw, ch = Inches(9), Inches(4.6)
    cx = (W - cw) / 2
    cy = (H - ch) / 2
    add_rect(sl, cx, cy, cw, ch, WHITE)
    add_rect(sl, cx, cy, cw, Inches(0.08), ACCENT)
    add_text(sl, "🌿  村上農園", cx, cy + Inches(0.3), cw, Inches(0.7),
             font_size=28, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "タスク管理システム", cx, cy + Inches(0.95), cw, Inches(0.7),
             font_size=32, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_rect(sl, cx + Inches(3.2), cy + Inches(1.85), Inches(2.6), Inches(0.05), LIGHT_GREEN)
    add_text(sl, "スタッフ向け  操作ガイド", cx, cy + Inches(1.95), cw, Inches(0.6),
             font_size=20, color=MED_GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, "スマートフォンからいつでも確認できます", cx, cy + Inches(2.7), cw, Inches(0.5),
             font_size=14, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(sl, "2026年6月", cx, cy + Inches(3.9), cw, Inches(0.4),
             font_size=12, color=SLATE, align=PP_ALIGN.CENTER)

    # ── スライド2：ログインとアクセス ──────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    header_bar(sl, "アクセスとログイン", "スマートフォンのブラウザから使います")

    # URL表示
    add_rect(sl, Inches(0.4), Inches(1.45), Inches(12.5), Inches(1.1), WHITE, ACCENT, 2)
    add_text(sl, "アクセス先URL", Inches(0.6), Inches(1.55), Inches(3), Inches(0.35),
             font_size=12, bold=True, color=ACCENT)
    add_text(sl, "https://kengo0127.github.io/Murakami_farm/",
             Inches(0.6), Inches(1.85), Inches(12), Inches(0.4),
             font_size=16, bold=True, color=MED_GREEN)

    # 左：手順
    section_badge(sl, "ログイン手順", Inches(0.4), Inches(2.75))
    steps = [
        ("①", "ブラウザでURLを開く\n（ブックマーク登録を推奨！）"),
        ("②", "メールアドレス（ユーザーID）と\nパスワードを入力する"),
        ("③", "「ログイン」ボタンを押す\n→ 当日のタスク一覧が表示される"),
    ]
    for i, (num, txt) in enumerate(steps):
        ty = Inches(3.2) + i * Inches(1.1)
        add_rect(sl, Inches(0.4), ty, Inches(0.55), Inches(0.55), MED_GREEN)
        add_text(sl, num, Inches(0.4), ty, Inches(0.55), Inches(0.55),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, txt, Inches(1.05), ty + Inches(0.04), Inches(5.2), Inches(0.9),
                 font_size=14, color=DARK_GRAY)

    # 右：ポイント
    add_rect(sl, Inches(7.0), Inches(2.75), Inches(5.9), Inches(4.35), WHITE, MED_GREEN, 1.5)
    add_text(sl, "📌  ご注意", Inches(7.2), Inches(2.85), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=MED_GREEN)
    notes = [
        "・ メールアドレスとパスワードは\n  管理者から配布されたものを使用",
        "・ パスワードは他の人に教えないでください",
        "・ 作業終了後は「ログアウト」を押してください",
        "・ ログイン情報がわからない場合は\n  管理者にご連絡ください",
    ]
    add_bullets(sl, notes, Inches(7.2), Inches(3.3), Inches(5.5), Inches(3.5), font_size=13)

    # ── スライド3：タスク確認と完了報告 ────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, LIGHT_GRAY)
    header_bar(sl, "タスクの確認と完了報告", "ログイン後の画面の使い方")

    # 左：タスク画面の説明
    section_badge(sl, "画面の見方", Inches(0.4), Inches(1.45))
    items = [
        "📅  今日の担当タスクが時刻の早い順に表示される",
        "🟢  緑の枠 ＝「▶ 作業中」（今やる作業）",
        "☕  黄色のカード ＝ 休憩時間",
        "✅  グレーのカード ＝ 完了済み",
        "⏭  次にやる作業は「次の作業」と表示される",
    ]
    add_bullets(sl, items, Inches(0.4), Inches(1.9), Inches(6.0), Inches(3.0), font_size=14)

    # 完了ボタン手順
    section_badge(sl, "完了の報告方法", Inches(0.4), Inches(5.0), bg=DARK_GREEN)
    comp_steps = [
        ("①", "作業が終わったら「✓ 完了」ボタンを押す"),
        ("②", "ボタンが ✅ 完了済み に変わる"),
        ("③", "次の作業が自動でハイライトされる"),
    ]
    for i, (num, txt) in enumerate(comp_steps):
        ty = Inches(5.45) + i * Inches(0.55)
        add_rect(sl, Inches(0.4), ty, Inches(0.4), Inches(0.4), ACCENT)
        add_text(sl, num, Inches(0.4), ty, Inches(0.4), Inches(0.4),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, txt, Inches(0.9), ty + Inches(0.04), Inches(5.4), Inches(0.38),
                 font_size=13, color=DARK_GRAY)

    # 右：重要注意事項
    add_rect(sl, Inches(7.0), Inches(1.45), Inches(5.9), Inches(5.65), WHITE, RED, 2)
    add_text(sl, "⚠️  重要なポイント", Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.45),
             font_size=15, bold=True, color=RED)
    warnings = [
        "・ 完了ボタンは1タスクに1回だけ",
        "・ チームの誰か1人が押せばOK\n  （全員の画面に反映される）",
        "・ 完了を取り消すことはできません\n  間違えた場合は管理者に連絡を",
        "・ 作業が終わっていなくても\n  他の人が押す場合があります",
    ]
    add_bullets(sl, warnings, Inches(7.2), Inches(2.15), Inches(5.5), Inches(4.0),
                font_size=13, color=DARK_GRAY)

    prs.save(path)
    print("staff ppt saved")


# ── メイン実行 ────────────────────────────────────────────────
import os
base = r"C:\Users\oshms\HiPro\02_村上農園\アプリ開発用"
make_admin_ppt(os.path.join(base, "村上農園_管理者向けガイド.pptx"))
make_staff_ppt(os.path.join(base, "村上農園_スタッフ向けガイド.pptx"))
print("done")
